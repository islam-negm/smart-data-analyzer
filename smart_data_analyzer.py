"""
Smart Data Analyzer - نظام ذكي متكامل لتحليل البيانات
================================================
يقوم بتحليل ملفات Excel تلقائياً وإنشاء:
- تقرير PDF احترافي
- ملف Excel بالنتائج
- عرض PowerPoint تقديمي
- Dashboard تفاعلي (Streamlit)
- تقرير تفسيري بالعربية
"""

import os
import sys
import json
import warnings
import textwrap
from datetime import datetime
from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.gridspec import GridSpec
from sklearn.linear_model import LinearRegression
from sklearn.preprocessing import StandardScaler
from sklearn.cluster import KMeans

warnings.filterwarnings('ignore')

# ─────────────────────────────────────────────
# إعداد المسارات
# ─────────────────────────────────────────────
BASE_DIR   = Path(__file__).parent
OUTPUT_DIR = BASE_DIR / "outputs"
ASSETS_DIR = BASE_DIR / "assets"
OUTPUT_DIR.mkdir(exist_ok=True)
ASSETS_DIR.mkdir(exist_ok=True)

TIMESTAMP = datetime.now().strftime("%Y%m%d_%H%M%S")

# ─────────────────────────────────────────────
# الألوان الاحترافية
# ─────────────────────────────────────────────
PALETTE = {
    "primary":   "#1B4F72",
    "secondary": "#2E86C1",
    "accent":    "#F39C12",
    "success":   "#27AE60",
    "danger":    "#E74C3C",
    "light":     "#EBF5FB",
    "dark":      "#1A252F",
    "gray":      "#7F8C8D",
    "white":     "#FFFFFF",
}
COLOR_LIST = ["#2E86C1","#F39C12","#27AE60","#E74C3C","#8E44AD",
              "#17A589","#D35400","#2C3E50","#85C1E9","#F8C471"]


# ══════════════════════════════════════════════
# MODULE 1 – قراءة وتنظيف البيانات
# ══════════════════════════════════════════════
class DataLoader:
    """تحميل وتنظيف ملفات Excel"""

    def __init__(self, filepath: str):
        self.filepath = Path(filepath)
        self.raw: dict[str, pd.DataFrame] = {}
        self.clean: dict[str, pd.DataFrame] = {}
        self.meta: dict = {}

    def load(self) -> "DataLoader":
        print(f"📂  جارٍ تحميل: {self.filepath.name}")
        xl = pd.ExcelFile(self.filepath)
        for sheet in xl.sheet_names:
            try:
                df = xl.parse(sheet)
                self.raw[sheet] = df
                print(f"    ✅  ورقة '{sheet}': {df.shape[0]} صف × {df.shape[1]} عمود")
            except Exception as e:
                print(f"    ⚠️  تعذّر قراءة '{sheet}': {e}")
        return self

    def clean_data(self) -> "DataLoader":
        for name, df in self.raw.items():
            df = df.copy()
            # حذف الأعمدة والصفوف الفارغة كلياً
            df.dropna(how="all", axis=0, inplace=True)
            df.dropna(how="all", axis=1, inplace=True)
            # تنظيف أسماء الأعمدة
            df.columns = [str(c).strip() for c in df.columns]
            # ملء القيم المفقودة الرقمية بالوسيط
            num_cols = df.select_dtypes(include=[np.number]).columns
            df[num_cols] = df[num_cols].fillna(df[num_cols].median())
            # ملء النصية بـ "غير محدد"
            cat_cols = df.select_dtypes(exclude=[np.number]).columns
            df[cat_cols] = df[cat_cols].fillna("غير محدد")
            # إزالة التكرارات
            before = len(df)
            df.drop_duplicates(inplace=True)
            after = len(df)
            self.clean[name] = df
            self.meta[name] = {
                "rows_raw":        before,
                "rows_clean":      after,
                "duplicates_removed": before - after,
                "missing_filled":  int(df.isnull().sum().sum()),
                "columns":         list(df.columns),
            }
        print(f"\n🧹  تنظيف البيانات: {len(self.clean)} ورقة جاهزة\n")
        return self


# ══════════════════════════════════════════════
# MODULE 2 – كشف نوع البيانات وتصنيف الأعمدة
# ══════════════════════════════════════════════
class ColumnClassifier:
    """يحدد تلقائياً أعمدة التاريخ، القيمة، الفئة، الكمية"""

    DATE_KEYWORDS  = ["date","تاريخ","month","شهر","year","سنة","quarter","ربع","period","فترة","week","أسبوع"]
    VALUE_KEYWORDS = ["sales","مبيعات","revenue","إيراد","amount","مبلغ","total","إجمالي","price","سعر",
                      "profit","ربح","cost","تكلفة","salary","راتب","income","دخل","budget","ميزانية",
                      "expense","مصروف","value","قيمة","score","درجة","rate","معدل"]
    QTY_KEYWORDS   = ["quantity","كمية","count","عدد","units","وحدات","orders","طلبات","transactions","معاملات",
                      "students","طلاب","employees","موظفين","customers","عملاء"]
    CAT_KEYWORDS   = ["category","فئة","region","منطقة","product","منتج","department","قسم",
                      "type","نوع","status","حالة","gender","جنس","grade","صف","level","مستوى"]

    def classify(self, df: pd.DataFrame) -> dict:
        cols = {"date": [], "value": [], "quantity": [], "category": [], "other": []}
        for col in df.columns:
            low = col.lower()
            dtype = df[col].dtype
            if dtype in ["datetime64[ns]"] or any(k in low for k in self.DATE_KEYWORDS):
                cols["date"].append(col)
            elif any(k in low for k in self.VALUE_KEYWORDS) and pd.api.types.is_numeric_dtype(dtype):
                cols["value"].append(col)
            elif any(k in low for k in self.QTY_KEYWORDS) and pd.api.types.is_numeric_dtype(dtype):
                cols["quantity"].append(col)
            elif any(k in low for k in self.CAT_KEYWORDS) or df[col].nunique() < 20:
                cols["category"].append(col)
            elif pd.api.types.is_numeric_dtype(dtype):
                cols["value"].append(col)
            else:
                cols["other"].append(col)
        return cols


# ══════════════════════════════════════════════
# MODULE 3 – التحليل الإحصائي الشامل
# ══════════════════════════════════════════════
class StatisticalAnalyzer:
    """يُنتج ملخصاً إحصائياً شاملاً لكل ورقة"""

    def analyze(self, df: pd.DataFrame, col_map: dict) -> dict:
        result = {}

        # ملخص أساسي
        num_df = df.select_dtypes(include=[np.number])
        result["summary"] = {
            "total_rows":     len(df),
            "total_columns":  len(df.columns),
            "numeric_columns":len(num_df.columns),
        }

        if num_df.empty:
            return result

        # إحصاء كل عمود قيمي
        val_stats = {}
        for col in (col_map["value"] + col_map["quantity"]):
            if col not in df.columns:
                continue
            s = df[col].dropna()
            if s.empty:
                continue
            val_stats[col] = {
                "total":  round(float(s.sum()), 2),
                "mean":   round(float(s.mean()), 2),
                "median": round(float(s.median()), 2),
                "std":    round(float(s.std()), 2),
                "min":    round(float(s.min()), 2),
                "max":    round(float(s.max()), 2),
                "q25":    round(float(s.quantile(.25)), 2),
                "q75":    round(float(s.quantile(.75)), 2),
            }
        result["column_stats"] = val_stats

        # أعلى فئات مساهمةً
        category_insights = {}
        for cat in col_map["category"]:
            if cat not in df.columns:
                continue
            for val in col_map["value"][:2]:
                if val not in df.columns:
                    continue
                grp = df.groupby(cat)[val].sum().sort_values(ascending=False)
                if grp.empty:
                    continue
                category_insights[f"{cat}_vs_{val}"] = {
                    "top":    grp.head(5).to_dict(),
                    "bottom": grp.tail(3).to_dict(),
                }
        result["category_insights"] = category_insights

        # اتجاه زمني
        if col_map["date"] and col_map["value"]:
            date_col = col_map["date"][0]
            val_col  = col_map["value"][0]
            try:
                tmp = df[[date_col, val_col]].copy()
                tmp[date_col] = pd.to_datetime(tmp[date_col], errors="coerce")
                tmp.dropna(subset=[date_col], inplace=True)
                tmp.sort_values(date_col, inplace=True)
                monthly = tmp.groupby(pd.Grouper(key=date_col, freq="ME"))[val_col].sum()
                if len(monthly) >= 2:
                    x = np.arange(len(monthly)).reshape(-1, 1)
                    y = monthly.values
                    lr = LinearRegression().fit(x, y)
                    trend = "تصاعدي 📈" if lr.coef_[0] > 0 else "تنازلي 📉"
                    result["trend"] = {
                        "direction": trend,
                        "slope":     round(float(lr.coef_[0]), 2),
                        "r2":        round(float(lr.score(x, y)), 3),
                        "monthly":   monthly.to_dict(),
                    }
            except Exception:
                pass

        # تجميع K-Means (إن وُجد بيانات كافية)
        if len(num_df) >= 10 and len(num_df.columns) >= 2:
            try:
                scaler = StandardScaler()
                scaled = scaler.fit_transform(num_df.fillna(0).iloc[:, :5])
                km = KMeans(n_clusters=min(3, len(num_df)), random_state=42, n_init=10)
                labels = km.fit_predict(scaled)
                result["clusters"] = {
                    "count":  int(len(set(labels))),
                    "sizes":  {str(i): int((labels == i).sum()) for i in set(labels)},
                }
            except Exception:
                pass

        return result


# ══════════════════════════════════════════════
# MODULE 4 – إنشاء الرسوم البيانية
# ══════════════════════════════════════════════
class ChartGenerator:
    """يُنشئ مجموعة متكاملة من الرسوم البيانية ويحفظها"""

    def __init__(self, output_dir: Path):
        self.output_dir = output_dir
        self.saved: list[str] = []

    def _save(self, fig: plt.Figure, name: str) -> str:
        path = str(self.output_dir / f"{name}_{TIMESTAMP}.png")
        fig.savefig(path, dpi=150, bbox_inches="tight", facecolor=PALETTE["white"])
        plt.close(fig)
        self.saved.append(path)
        return path

    # ── Bar Chart ──────────────────────────────
    def bar_chart(self, df: pd.DataFrame, cat_col: str, val_col: str, title: str) -> str:
        grp = df.groupby(cat_col)[val_col].sum().sort_values(ascending=False).head(10)
        fig, ax = plt.subplots(figsize=(10, 6))
        bars = ax.bar(range(len(grp)), grp.values,
                      color=[COLOR_LIST[i % len(COLOR_LIST)] for i in range(len(grp))],
                      edgecolor="white", linewidth=0.8)
        ax.set_xticks(range(len(grp)))
        ax.set_xticklabels([str(x)[:18] for x in grp.index], rotation=30, ha="right", fontsize=9)
        ax.set_title(title, fontsize=14, fontweight="bold", color=PALETTE["dark"], pad=15)
        ax.set_ylabel(val_col, fontsize=10)
        ax.spines[["top","right"]].set_visible(False)
        ax.set_facecolor(PALETTE["light"])
        # قيمة فوق كل عمود
        for bar, v in zip(bars, grp.values):
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() * 1.01,
                    f"{v:,.0f}", ha="center", va="bottom", fontsize=8, color=PALETTE["dark"])
        fig.tight_layout()
        return self._save(fig, "bar_chart")

    # ── Line Chart ─────────────────────────────
    def line_chart(self, monthly: dict, val_col: str, title: str) -> str:
        dates  = list(monthly.keys())
        values = list(monthly.values())
        fig, ax = plt.subplots(figsize=(11, 5))
        ax.fill_between(range(len(dates)), values, alpha=0.15, color=PALETTE["secondary"])
        ax.plot(range(len(dates)), values, marker="o", color=PALETTE["secondary"],
                linewidth=2.5, markersize=6)
        step = max(1, len(dates) // 10)
        ax.set_xticks(range(0, len(dates), step))
        ax.set_xticklabels([str(d)[:7] for d in dates[::step]], rotation=30, ha="right", fontsize=8)
        ax.set_title(title, fontsize=14, fontweight="bold", color=PALETTE["dark"], pad=15)
        ax.set_ylabel(val_col, fontsize=10)
        ax.spines[["top","right"]].set_visible(False)
        ax.set_facecolor(PALETTE["light"])
        fig.tight_layout()
        return self._save(fig, "line_chart")

    # ── Pie Chart ──────────────────────────────
    def pie_chart(self, df: pd.DataFrame, cat_col: str, val_col: str, title: str) -> str:
        grp = df.groupby(cat_col)[val_col].sum().sort_values(ascending=False).head(8)
        fig, ax = plt.subplots(figsize=(8, 8))
        wedges, texts, autotexts = ax.pie(
            grp.values,
            labels=[str(x)[:20] for x in grp.index],
            autopct="%1.1f%%",
            colors=COLOR_LIST[:len(grp)],
            pctdistance=0.8,
            startangle=140,
            wedgeprops=dict(edgecolor="white", linewidth=1.5)
        )
        for t in autotexts:
            t.set_fontsize(9)
            t.set_color("white")
            t.set_fontweight("bold")
        ax.set_title(title, fontsize=14, fontweight="bold", color=PALETTE["dark"], pad=20)
        fig.tight_layout()
        return self._save(fig, "pie_chart")

    # ── Scatter Plot ───────────────────────────
    def scatter_plot(self, df: pd.DataFrame, x_col: str, y_col: str, title: str) -> str:
        fig, ax = plt.subplots(figsize=(9, 6))
        ax.scatter(df[x_col], df[y_col], alpha=0.6, s=60,
                   color=PALETTE["secondary"], edgecolors=PALETTE["primary"], linewidth=0.5)
        # خط اتجاه
        try:
            z = np.polyfit(df[x_col].dropna(), df[y_col].dropna(), 1)
            p = np.poly1d(z)
            xs = np.linspace(df[x_col].min(), df[x_col].max(), 100)
            ax.plot(xs, p(xs), "--", color=PALETTE["danger"], linewidth=1.5, label="Trend")
        except Exception:
            pass
        ax.set_xlabel(x_col, fontsize=10)
        ax.set_ylabel(y_col, fontsize=10)
        ax.set_title(title, fontsize=14, fontweight="bold", color=PALETTE["dark"], pad=15)
        ax.spines[["top","right"]].set_visible(False)
        ax.set_facecolor(PALETTE["light"])
        fig.tight_layout()
        return self._save(fig, "scatter_plot")

    # ── Heatmap ────────────────────────────────
    def heatmap(self, df: pd.DataFrame, title: str) -> str:
        num_df = df.select_dtypes(include=[np.number]).dropna(axis=1)
        if num_df.shape[1] < 2:
            return ""
        corr = num_df.corr()
        fig, ax = plt.subplots(figsize=(max(8, len(corr)), max(6, len(corr) - 1)))
        im = ax.imshow(corr.values, cmap="RdYlGn", vmin=-1, vmax=1, aspect="auto")
        fig.colorbar(im, ax=ax, shrink=0.8)
        ax.set_xticks(range(len(corr.columns)))
        ax.set_yticks(range(len(corr.columns)))
        ax.set_xticklabels([c[:15] for c in corr.columns], rotation=45, ha="right", fontsize=8)
        ax.set_yticklabels([c[:15] for c in corr.columns], fontsize=8)
        for i in range(len(corr)):
            for j in range(len(corr)):
                v = corr.values[i, j]
                ax.text(j, i, f"{v:.2f}", ha="center", va="center",
                        fontsize=7, color="black" if abs(v) < 0.6 else "white")
        ax.set_title(title, fontsize=13, fontweight="bold", color=PALETTE["dark"], pad=15)
        fig.tight_layout()
        return self._save(fig, "heatmap")

    # ── Dashboard ──────────────────────────────
    def dashboard(self, df: pd.DataFrame, col_map: dict, analysis: dict, sheet_name: str) -> str:
        fig = plt.figure(figsize=(18, 14), facecolor=PALETTE["dark"])
        fig.suptitle(f"📊  لوحة معلومات – {sheet_name}",
                     fontsize=20, fontweight="bold", color=PALETTE["white"], y=0.98)
        gs = GridSpec(3, 3, figure=fig, hspace=0.45, wspace=0.35)

        # KPIs شريط علوي
        kpi_ax = fig.add_subplot(gs[0, :])
        kpi_ax.set_facecolor(PALETTE["dark"])
        kpi_ax.axis("off")
        kpis = []
        for col, stats in list(analysis.get("column_stats", {}).items())[:4]:
            kpis.append((col, stats.get("total", 0), stats.get("mean", 0)))
        xpos = [0.1, 0.35, 0.6, 0.85]
        for i, (col, total, mean) in enumerate(kpis[:4]):
            kpi_ax.text(xpos[i], 0.75, f"{total:,.0f}",
                        transform=kpi_ax.transAxes, fontsize=22, fontweight="bold",
                        color=COLOR_LIST[i], ha="center")
            kpi_ax.text(xpos[i], 0.3, col[:20],
                        transform=kpi_ax.transAxes, fontsize=9,
                        color=PALETTE["gray"], ha="center")

        # Bar Chart (أيسر وسط)
        ax1 = fig.add_subplot(gs[1, 0])
        if col_map["category"] and col_map["value"]:
            c, v = col_map["category"][0], col_map["value"][0]
            grp = df.groupby(c)[v].sum().sort_values(ascending=False).head(6)
            ax1.barh(range(len(grp)), grp.values, color=COLOR_LIST[:len(grp)])
            ax1.set_yticks(range(len(grp)))
            ax1.set_yticklabels([str(x)[:12] for x in grp.index], fontsize=8, color="white")
            ax1.set_facecolor(PALETTE["dark"])
            ax1.tick_params(colors="white")
            ax1.spines[:].set_color(PALETTE["gray"])
        ax1.set_title("Top Categories", fontsize=10, color=PALETTE["white"])

        # Pie Chart (وسط)
        ax2 = fig.add_subplot(gs[1, 1])
        if col_map["category"] and col_map["value"]:
            c, v = col_map["category"][0], col_map["value"][0]
            grp = df.groupby(c)[v].sum().sort_values(ascending=False).head(5)
            ax2.pie(grp.values, labels=[str(x)[:10] for x in grp.index],
                    colors=COLOR_LIST[:len(grp)], autopct="%1.0f%%",
                    pctdistance=0.8, wedgeprops=dict(edgecolor=PALETTE["dark"], linewidth=1.5),
                    textprops=dict(color="white", fontsize=8))
        ax2.set_facecolor(PALETTE["dark"])
        ax2.set_title("Distribution", fontsize=10, color=PALETTE["white"])

        # Line Chart (أيمن)
        ax3 = fig.add_subplot(gs[1, 2])
        if "trend" in analysis and "monthly" in analysis["trend"]:
            monthly = analysis["trend"]["monthly"]
            vals = list(monthly.values())
            ax3.plot(range(len(vals)), vals, color=PALETTE["accent"],
                     linewidth=2, marker="o", markersize=4)
            ax3.fill_between(range(len(vals)), vals, alpha=0.15, color=PALETTE["accent"])
            ax3.set_facecolor(PALETTE["dark"])
            ax3.tick_params(colors="white")
            ax3.spines[:].set_color(PALETTE["gray"])
            ax3.set_title(f"الاتجاه {analysis['trend']['direction']}", fontsize=10, color=PALETTE["white"])

        # Heatmap (أسفل)
        ax4 = fig.add_subplot(gs[2, :])
        num_df = df.select_dtypes(include=[np.number]).dropna(axis=1).iloc[:, :8]
        if num_df.shape[1] >= 2:
            corr = num_df.corr()
            im = ax4.imshow(corr.values, cmap="coolwarm", vmin=-1, vmax=1, aspect="auto")
            ax4.set_xticks(range(len(corr.columns)))
            ax4.set_yticks(range(len(corr.columns)))
            ax4.set_xticklabels([c[:12] for c in corr.columns],
                                 rotation=30, ha="right", fontsize=8, color="white")
            ax4.set_yticklabels([c[:12] for c in corr.columns], fontsize=8, color="white")
            for i in range(len(corr)):
                for j in range(len(corr)):
                    ax4.text(j, i, f"{corr.values[i,j]:.2f}",
                             ha="center", va="center", fontsize=7,
                             color="white" if abs(corr.values[i,j]) > 0.5 else "black")
            fig.colorbar(im, ax=ax4, orientation="horizontal", shrink=0.4, pad=0.05)
        ax4.set_title("مصفوفة الارتباط", fontsize=10, color=PALETTE["white"])
        ax4.set_facecolor(PALETTE["dark"])

        return self._save(fig, "dashboard")


# ══════════════════════════════════════════════
# MODULE 5 – التقرير التفسيري بالعربية
# ══════════════════════════════════════════════
class ArabicReportWriter:
    """يكتب تقريراً تفسيرياً شاملاً بالعربية"""

    def write(self, loader: DataLoader, all_analysis: dict) -> str:
        now = datetime.now().strftime("%Y/%m/%d  %H:%M")
        lines = []

        def h(text, level=1):
            sep = "═" * 60 if level == 1 else "─" * 50
            lines.append(sep)
            lines.append(f"{'  ' * (level-1)}{text}")
            lines.append(sep)

        def p(text):
            lines.append(text)

        def nl():
            lines.append("")

        h("📊  تقرير تحليل البيانات الذكي – Smart Data Analyzer")
        p(f"📅  تاريخ التقرير: {now}")
        p(f"📁  الملف: {loader.filepath.name}")
        nl()

        # ملخص تنفيذي
        h("ملخص تنفيذي", 2)
        total_rows = sum(m["rows_clean"] for m in loader.meta.values())
        total_cols = sum(len(m["columns"]) for m in loader.meta.values())
        total_dup  = sum(m["duplicates_removed"] for m in loader.meta.values())
        p(f"• إجمالي الأوراق المُحللة  : {len(loader.clean)}")
        p(f"• إجمالي السجلات المنظّفة  : {total_rows:,}")
        p(f"• إجمالي الأعمدة           : {total_cols}")
        p(f"• السجلات المكررة المحذوفة : {total_dup:,}")
        nl()

        for sheet, analysis in all_analysis.items():
            h(f"📋  ورقة: {sheet}", 2)
            meta = loader.meta.get(sheet, {})
            p(f"• الصفوف بعد التنظيف : {meta.get('rows_clean', '-'):,}")
            p(f"• الأعمدة             : {len(meta.get('columns', []))}")
            nl()

            # إحصاء الأعمدة
            col_stats = analysis.get("column_stats", {})
            if col_stats:
                h(f"المؤشرات الرقمية الرئيسية", 3)
                for col, stats in col_stats.items():
                    p(f"\n  ► {col}")
                    p(f"    - الإجمالي  : {stats['total']:>15,.2f}")
                    p(f"    - المتوسط   : {stats['mean']:>15,.2f}")
                    p(f"    - الوسيط    : {stats['median']:>15,.2f}")
                    p(f"    - الانحراف  : {stats['std']:>15,.2f}")
                    p(f"    - الأدنى    : {stats['min']:>15,.2f}")
                    p(f"    - الأعلى    : {stats['max']:>15,.2f}")
                nl()

            # تحليل الفئات
            cat_insights = analysis.get("category_insights", {})
            if cat_insights:
                h("تحليل الفئات والتوزيع", 3)
                for key, insight in cat_insights.items():
                    parts = key.split("_vs_")
                    cat_name = parts[0] if parts else key
                    val_name = parts[1] if len(parts) > 1 else ""
                    p(f"\n  ► أعلى {cat_name} من حيث {val_name}:")
                    for k, v in list(insight["top"].items())[:5]:
                        bar_len = int((v / max(insight["top"].values())) * 20) if insight["top"].values() else 0
                        p(f"    {'█' * bar_len:<20}  {str(k)[:25]:<25}  {v:>12,.0f}")
                nl()

            # الاتجاه الزمني
            trend = analysis.get("trend", {})
            if trend:
                h("تحليل الاتجاه الزمني", 3)
                p(f"• الاتجاه العام   : {trend['direction']}")
                p(f"• معدل الانحدار   : {trend['slope']:,.2f}  (وحدة لكل فترة)")
                p(f"• معامل التحديد R²: {trend['r2']:.3f}  "
                  f"({'قوي' if trend['r2'] > 0.7 else 'متوسط' if trend['r2'] > 0.4 else 'ضعيف'})")
                nl()

            # التجميعات
            clusters = analysis.get("clusters", {})
            if clusters:
                h("تحليل التجميعات (K-Means)", 3)
                p(f"• عدد المجموعات المكتشفة : {clusters['count']}")
                for cl, sz in clusters["sizes"].items():
                    p(f"    - المجموعة {int(cl)+1} : {sz:,} سجل  ({sz/meta.get('rows_clean',1)*100:.1f}%)")
                nl()

            # توصيات تلقائية
            h("التوصيات والاستنتاجات", 3)
            if trend:
                if "تصاعدي" in trend["direction"]:
                    p("✅  البيانات تُظهر نمواً إيجابياً مستمراً؛ يُنصح بالتوسع والاستثمار.")
                else:
                    p("⚠️  البيانات تُظهر انخفاضاً؛ يُنصح بمراجعة الاستراتيجية الحالية.")
            if col_stats:
                high_std = [(c, s["std"]/s["mean"] if s["mean"] else 0)
                            for c, s in col_stats.items() if s.get("mean", 0) > 0]
                high_std.sort(key=lambda x: x[1], reverse=True)
                if high_std and high_std[0][1] > 0.5:
                    p(f"⚠️  عمود '{high_std[0][0]}' يعاني من تشتت عالٍ (CV={high_std[0][1]:.2f})؛ "
                      f"يُنصح بالتعمق في أسباب التفاوت.")
            p("📌  يُوصى بمراجعة هذا التقرير مع الفريق المعني واتخاذ القرارات بناءً على النتائج.")
            nl()

        h("نهاية التقرير", 2)
        p("تم إنشاء هذا التقرير تلقائياً بواسطة نظام Smart Data Analyzer")
        p(f"© {datetime.now().year}  جميع الحقوق محفوظة")

        return "\n".join(lines)


# ══════════════════════════════════════════════
# MODULE 6 – إنشاء تقرير PDF
# ══════════════════════════════════════════════
class PDFReportGenerator:
    """يُنشئ تقرير PDF احترافي بالعربية باستخدام reportlab"""

    def __init__(self, output_dir: Path):
        self.output_dir = output_dir

    def generate(self, report_text: str, charts: list[str], sheet_name: str) -> str:
        from reportlab.lib.pagesizes    import A4
        from reportlab.lib.units        import cm
        from reportlab.lib              import colors
        from reportlab.lib.styles       import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus         import (SimpleDocTemplate, Paragraph, Spacer,
                                                 Image, HRFlowable, Table, TableStyle,
                                                 PageBreak)
        from reportlab.lib.enums        import TA_CENTER, TA_LEFT, TA_RIGHT
        from reportlab.pdfbase          import pdfmetrics
        from reportlab.pdfbase.ttfonts  import TTFont

        out_path = str(self.output_dir / f"report_{TIMESTAMP}.pdf")
        doc = SimpleDocTemplate(out_path, pagesize=A4,
                                leftMargin=2*cm, rightMargin=2*cm,
                                topMargin=2.5*cm, bottomMargin=2.5*cm)

        # ── الأنماط ──
        styles = getSampleStyleSheet()

        def style(name, **kw):
            return ParagraphStyle(name, parent=styles["Normal"], **kw)

        title_s  = style("MyTitle",  fontSize=22, alignment=TA_CENTER,
                          textColor=colors.HexColor(PALETTE["primary"]),
                          fontName="Helvetica-Bold", spaceAfter=6)
        h1_s     = style("MyH1",     fontSize=14, alignment=TA_RIGHT,
                          textColor=colors.HexColor(PALETTE["primary"]),
                          fontName="Helvetica-Bold", spaceBefore=12, spaceAfter=4)
        h2_s     = style("MyH2",     fontSize=11, alignment=TA_RIGHT,
                          textColor=colors.HexColor(PALETTE["secondary"]),
                          fontName="Helvetica-Bold", spaceBefore=8, spaceAfter=3)
        body_s   = style("MyBody",   fontSize=9,  alignment=TA_RIGHT,
                          textColor=colors.HexColor(PALETTE["dark"]),
                          leading=16, spaceAfter=2)
        caption_s= style("MyCaption",fontSize=8,  alignment=TA_CENTER,
                          textColor=colors.HexColor(PALETTE["gray"]),
                          fontName="Helvetica-Oblique")

        story = []

        # ── غلاف ──
        story.append(Spacer(1, 1.5*cm))
        story.append(Paragraph("Smart Data Analyzer", title_s))
        story.append(Paragraph("تقرير تحليل البيانات الذكي", title_s))
        story.append(HRFlowable(width="100%", thickness=2,
                                color=colors.HexColor(PALETTE["accent"])))
        story.append(Spacer(1, 0.3*cm))
        story.append(Paragraph(f"تاريخ الإنشاء: {datetime.now().strftime('%Y/%m/%d')}", caption_s))
        story.append(Spacer(1, 1*cm))

        # ── النص التفسيري ──
        story.append(Paragraph("التقرير التفسيري الشامل", h1_s))
        story.append(HRFlowable(width="80%", thickness=1,
                                color=colors.HexColor(PALETTE["secondary"])))
        story.append(Spacer(1, 0.3*cm))

        for line in report_text.split("\n"):
            line = line.strip()
            if not line:
                story.append(Spacer(1, 0.15*cm))
                continue
            if line.startswith("═"):
                story.append(HRFlowable(width="100%", thickness=1.5,
                                        color=colors.HexColor(PALETTE["primary"])))
            elif line.startswith("─"):
                story.append(HRFlowable(width="80%", thickness=0.5,
                                        color=colors.HexColor(PALETTE["gray"])))
            elif any(line.startswith(h) for h in ["📊","📋","📅","📁","✅","⚠️","📌","©"]):
                story.append(Paragraph(line, h2_s))
            elif line.startswith("•"):
                story.append(Paragraph(line, body_s))
            elif line.startswith("►"):
                story.append(Paragraph(f"<b>{line}</b>", body_s))
            elif line.startswith("-"):
                story.append(Paragraph(f"&nbsp;&nbsp;&nbsp;{line}", body_s))
            else:
                story.append(Paragraph(line, body_s))

        # ── الرسوم البيانية ──
        if charts:
            story.append(PageBreak())
            story.append(Paragraph("الرسوم البيانية والمخططات", h1_s))
            story.append(HRFlowable(width="80%", thickness=1,
                                    color=colors.HexColor(PALETTE["secondary"])))
            story.append(Spacer(1, 0.5*cm))

            chart_names = {
                "bar_chart":   "مخطط الأعمدة – توزيع الفئات",
                "line_chart":  "المخطط الخطي – الاتجاه الزمني",
                "pie_chart":   "المخطط الدائري – نسب المساهمة",
                "scatter_plot":"مخطط التشتت – العلاقة بين المتغيرات",
                "heatmap":     "خريطة الحرارة – مصفوفة الارتباط",
                "dashboard":   "لوحة المعلومات الشاملة",
            }
            for chart_path in charts:
                if not os.path.exists(chart_path):
                    continue
                chart_key = next((k for k in chart_names if k in chart_path), "")
                chart_title = chart_names.get(chart_key, "رسم بياني")
                try:
                    from PIL import Image as PILImage
                    pil_img = PILImage.open(chart_path)
                    w, h = pil_img.size
                    max_w = 15 * cm
                    ratio = max_w / w
                    img_h = h * ratio
                    img = Image(chart_path, width=max_w, height=img_h)
                    story.append(Paragraph(chart_title, h2_s))
                    story.append(img)
                    story.append(Spacer(1, 0.5*cm))
                except Exception:
                    pass

        doc.build(story)
        print(f"✅  PDF محفوظ: {out_path}")
        return out_path


# ══════════════════════════════════════════════
# MODULE 7 – إنشاء ملف Excel بالنتائج
# ══════════════════════════════════════════════
class ExcelResultsGenerator:
    """يُنشئ ملف Excel احترافي بجميع النتائج"""

    def __init__(self, output_dir: Path):
        self.output_dir = output_dir

    def generate(self, loader: DataLoader, all_analysis: dict) -> str:
        from openpyxl import Workbook
        from openpyxl.styles import (Font, PatternFill, Alignment,
                                      Border, Side, numbers)
        from openpyxl.utils import get_column_letter
        from openpyxl.chart import BarChart, Reference

        out_path = str(self.output_dir / f"results_{TIMESTAMP}.xlsx")
        wb = Workbook()
        wb.remove(wb.active)  # حذف الورقة الافتراضية

        def hdr_style(cell, color=PALETTE["primary"]):
            cell.font = Font(bold=True, color="FFFFFF", size=11)
            cell.fill = PatternFill("solid", fgColor=color.lstrip("#"))
            cell.alignment = Alignment(horizontal="center", vertical="center")

        def thin_border():
            s = Side(style="thin", color="CCCCCC")
            return Border(left=s, right=s, top=s, bottom=s)

        # ── ورقة الملخص التنفيذي ──
        ws = wb.create_sheet("الملخص التنفيذي")
        ws.sheet_view.rightToLeft = True
        ws.column_dimensions["A"].width = 35
        ws.column_dimensions["B"].width = 20

        ws.merge_cells("A1:B1")
        ws["A1"] = "ملخص تنفيذي – Smart Data Analyzer"
        hdr_style(ws["A1"])
        ws.row_dimensions[1].height = 30

        rows = [
            ("الملف المُحلَّل",        loader.filepath.name),
            ("تاريخ التحليل",          datetime.now().strftime("%Y/%m/%d %H:%M")),
            ("عدد الأوراق",            len(loader.clean)),
            ("إجمالي السجلات",         sum(m["rows_clean"] for m in loader.meta.values())),
            ("التكرارات المحذوفة",      sum(m["duplicates_removed"] for m in loader.meta.values())),
        ]
        for r, (lbl, val) in enumerate(rows, start=2):
            ws[f"A{r}"] = lbl
            ws[f"B{r}"] = val
            ws[f"A{r}"].font = Font(bold=True)
            ws[f"A{r}"].border = thin_border()
            ws[f"B{r}"].border = thin_border()
            ws[f"B{r}"].alignment = Alignment(horizontal="center")

        # ── ورقة لكل ورقة بيانات ──
        for sheet_name, df in loader.clean.items():
            # البيانات المنظّفة
            ws_data = wb.create_sheet(f"بيانات – {sheet_name}"[:30])
            ws_data.sheet_view.rightToLeft = True
            for ci, col in enumerate(df.columns, start=1):
                cell = ws_data.cell(row=1, column=ci, value=col)
                hdr_style(cell, PALETTE["secondary"])
                ws_data.column_dimensions[get_column_letter(ci)].width = 18
            for ri, row in enumerate(df.itertuples(index=False), start=2):
                for ci, val in enumerate(row, start=1):
                    c = ws_data.cell(row=ri, column=ci, value=val)
                    c.border = thin_border()
                    if ri % 2 == 0:
                        c.fill = PatternFill("solid", fgColor="EBF5FB")

            # الإحصاء
            analysis = all_analysis.get(sheet_name, {})
            ws_stat = wb.create_sheet(f"إحصاء – {sheet_name}"[:30])
            ws_stat.sheet_view.rightToLeft = True
            ws_stat.column_dimensions["A"].width = 22
            for ci, h in enumerate(["الحقل","الإجمالي","المتوسط","الوسيط","الانحراف","الأدنى","الأعلى"], start=1):
                cell = ws_stat.cell(row=1, column=ci, value=h)
                hdr_style(cell)
                ws_stat.column_dimensions[get_column_letter(ci)].width = 16
            for ri, (col, stats) in enumerate(analysis.get("column_stats", {}).items(), start=2):
                for ci, val in enumerate(
                    [col, stats["total"], stats["mean"], stats["median"],
                     stats["std"], stats["min"], stats["max"]], start=1
                ):
                    c = ws_stat.cell(row=ri, column=ci, value=val)
                    c.border = thin_border()
                    if ri % 2 == 0:
                        c.fill = PatternFill("solid", fgColor="EBF5FB")
                    if ci > 1:
                        c.number_format = "#,##0.00"

        wb.save(out_path)
        print(f"✅  Excel محفوظ: {out_path}")
        return out_path


# ══════════════════════════════════════════════
# MODULE 8 – عرض PowerPoint تقديمي
# ══════════════════════════════════════════════
class PowerPointGenerator:
    """يُنشئ عرضاً تقديمياً PowerPoint احترافياً"""

    def __init__(self, output_dir: Path):
        self.output_dir = output_dir

    def generate(self, loader: DataLoader, all_analysis: dict, charts: list[str]) -> str:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        def rgb(hex_color: str):
            h = hex_color.lstrip("#")
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]  # فارغ

        # ── دالة مساعدة لإضافة مستطيل ملوّن ──
        def add_rect(slide, l, t, w, h, fill_hex, alpha=None):
            shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(fill_hex)
            shape.line.fill.background()
            return shape

        # ── دالة إضافة نص ──
        def add_text(slide, text, l, t, w, h,
                     font_size=18, bold=False, color="#FFFFFF",
                     align=PP_ALIGN.LEFT, wrap=True):
            txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf  = txb.text_frame
            tf.word_wrap = wrap
            para = tf.paragraphs[0]
            para.alignment = align
            run = para.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
            return txb

        # ════ شريحة 1: الغلاف ════
        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, 13.33, 7.5, PALETTE["dark"])
        add_rect(slide, 0, 0, 13.33, 0.08, PALETTE["accent"])
        add_rect(slide, 0, 7.42, 13.33, 0.08, PALETTE["accent"])
        add_text(slide, "📊", 5.5, 0.8, 2, 1.2, font_size=54)
        add_text(slide, "Smart Data Analyzer", 1.5, 2.0, 10, 1.2,
                 font_size=42, bold=True, color=PALETTE["white"], align=PP_ALIGN.CENTER)
        add_text(slide, "نظام ذكي متكامل لتحليل البيانات وإنشاء التقارير",
                 1.5, 3.3, 10, 0.8, font_size=20, color=PALETTE["accent"], align=PP_ALIGN.CENTER)
        add_text(slide, f"الملف: {loader.filepath.name}  |  {datetime.now().strftime('%Y/%m/%d')}",
                 1.5, 4.3, 10, 0.6, font_size=14, color=PALETTE["gray"], align=PP_ALIGN.CENTER)

        # ════ شريحة 2: الملخص التنفيذي ════
        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, 13.33, 7.5, PALETTE["light"])
        add_rect(slide, 0, 0, 13.33, 1.3, PALETTE["primary"])
        add_text(slide, "الملخص التنفيذي", 0.5, 0.25, 12, 0.8,
                 font_size=28, bold=True, color=PALETTE["white"], align=PP_ALIGN.RIGHT)

        total_rows = sum(m["rows_clean"] for m in loader.meta.values())
        total_dup  = sum(m["duplicates_removed"] for m in loader.meta.values())

        kpi_data = [
            ("📋", str(len(loader.clean)), "عدد الأوراق"),
            ("📊", f"{total_rows:,}", "إجمالي السجلات"),
            ("🧹", f"{total_dup:,}", "تكرارات محذوفة"),
            ("📈", str(len(all_analysis)), "تحليلات منجزة"),
        ]
        for i, (icon, val, lbl) in enumerate(kpi_data):
            x = 0.5 + i * 3.2
            add_rect(slide, x, 1.6, 2.8, 2.0, PALETTE["primary"])
            add_text(slide, icon,  x+0.3, 1.7, 2, 0.8, font_size=28)
            add_text(slide, val,   x+0.1, 2.3, 2.5, 0.8,
                     font_size=26, bold=True, color=PALETTE["accent"], align=PP_ALIGN.CENTER)
            add_text(slide, lbl,   x+0.1, 3.1, 2.5, 0.5,
                     font_size=13, color=PALETTE["white"], align=PP_ALIGN.CENTER)

        # ════ شرائح التحليل لكل ورقة ════
        for sheet_name, analysis in all_analysis.items():
            slide = prs.slides.add_slide(blank)
            add_rect(slide, 0, 0, 13.33, 7.5, PALETTE["white"])
            add_rect(slide, 0, 0, 13.33, 1.2, PALETTE["secondary"])
            add_text(slide, f"📋  تحليل: {sheet_name}", 0.3, 0.2, 12, 0.8,
                     font_size=24, bold=True, color=PALETTE["white"], align=PP_ALIGN.RIGHT)

            # جدول الإحصاء
            col_stats = analysis.get("column_stats", {})
            y = 1.5
            for col, stats in list(col_stats.items())[:4]:
                add_rect(slide, 0.3, y, 12.5, 0.55, PALETTE["light"])
                add_text(slide, f"► {col}",       0.5, y+0.05, 3.5, 0.45,
                         font_size=11, bold=True, color=PALETTE["dark"])
                add_text(slide, f"الإجمالي: {stats['total']:,.0f}",  4.0, y+0.05, 2.5, 0.45,
                         font_size=11, color=PALETTE["primary"])
                add_text(slide, f"المتوسط: {stats['mean']:,.0f}",    6.5, y+0.05, 2.5, 0.45,
                         font_size=11, color=PALETTE["secondary"])
                add_text(slide, f"الأعلى: {stats['max']:,.0f}",      9.0, y+0.05, 2.5, 0.45,
                         font_size=11, color=PALETTE["success"])
                y += 0.65

            # الاتجاه
            trend = analysis.get("trend", {})
            if trend:
                add_rect(slide, 0.3, y+0.2, 12.5, 0.7, PALETTE["primary"])
                add_text(slide,
                         f"الاتجاه: {trend['direction']}   |   "
                         f"الانحدار: {trend['slope']:,.2f}   |   "
                         f"R² = {trend['r2']:.3f}",
                         0.5, y+0.25, 12, 0.6,
                         font_size=13, bold=True, color=PALETTE["accent"])

        # ════ شرائح الرسوم البيانية ════
        chart_titles = {
            "bar_chart":    "مخطط الأعمدة",
            "line_chart":   "المخطط الخطي",
            "pie_chart":    "المخطط الدائري",
            "scatter_plot": "مخطط التشتت",
            "heatmap":      "خريطة الحرارة",
            "dashboard":    "لوحة المعلومات",
        }
        for chart_path in charts:
            if not os.path.exists(chart_path):
                continue
            chart_key   = next((k for k in chart_titles if k in chart_path), "")
            chart_title = chart_titles.get(chart_key, "رسم بياني")
            slide = prs.slides.add_slide(blank)
            add_rect(slide, 0, 0, 13.33, 7.5, PALETTE["dark"])
            add_rect(slide, 0, 0, 13.33, 1.0, PALETTE["primary"])
            add_text(slide, chart_title, 0.5, 0.1, 12, 0.8,
                     font_size=24, bold=True, color=PALETTE["white"], align=PP_ALIGN.CENTER)
            try:
                slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.2),
                                         Inches(12.33), Inches(5.8))
            except Exception:
                pass

        # ════ شريحة الختام ════
        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, 13.33, 7.5, PALETTE["dark"])
        add_rect(slide, 0, 3.5, 13.33, 0.05, PALETTE["accent"])
        add_text(slide, "شكراً لاستخدام", 1.5, 1.5, 10, 1,
                 font_size=24, color=PALETTE["gray"], align=PP_ALIGN.CENTER)
        add_text(slide, "Smart Data Analyzer", 1.5, 2.5, 10, 1.2,
                 font_size=40, bold=True, color=PALETTE["accent"], align=PP_ALIGN.CENTER)
        add_text(slide, "نظام ذكي متكامل – من البيانات الخام إلى قرارات عملية",
                 1.5, 4.5, 10, 0.8, font_size=16,
                 color=PALETTE["white"], align=PP_ALIGN.CENTER)

        out_path = str(self.output_dir / f"presentation_{TIMESTAMP}.pptx")
        prs.save(out_path)
        print(f"✅  PowerPoint محفوظ: {out_path}")
        return out_path


# ══════════════════════════════════════════════
# MODULE 9 – إنشاء تطبيق Streamlit
# ══════════════════════════════════════════════
class StreamlitDashboardGenerator:
    """يُنشئ كود تطبيق Streamlit Dashboard"""

    def __init__(self, output_dir: Path):
        self.output_dir = output_dir

    def generate(self, loader_path: str) -> str:
        code = textwrap.dedent(f'''
        """
        Smart Data Analyzer – Dashboard تفاعلي
        تشغيل: streamlit run streamlit_dashboard.py
        """
        import streamlit as st
        import pandas as pd
        import numpy as np
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from pathlib import Path
        import sys, os

        st.set_page_config(
            page_title="Smart Data Analyzer",
            page_icon="📊",
            layout="wide",
            initial_sidebar_state="expanded",
        )

        # ── CSS مخصص ──
        st.markdown("""
        <style>
        .main {{background-color: #F0F4F8;}}
        .stMetric label {{font-size: 12px; color: #7F8C8D;}}
        .stMetric value {{font-size: 24px; font-weight: bold; color: #1B4F72;}}
        h1 {{color: #1B4F72; border-bottom: 3px solid #F39C12; padding-bottom: 8px;}}
        h2 {{color: #2E86C1;}}
        </style>
        """, unsafe_allow_html=True)

        # ── الشريط الجانبي ──
        with st.sidebar:
            st.image("https://via.placeholder.com/200x60/1B4F72/FFFFFF?text=Smart+Analyzer",
                     use_column_width=True)
            st.title("📁 رفع الملف")
            uploaded = st.file_uploader("اختر ملف Excel", type=["xlsx","xls"])
            st.markdown("---")
            st.info("📊 Smart Data Analyzer\\nنظام تحليل بيانات متكامل")

        st.title("📊 Smart Data Analyzer – لوحة المعلومات التفاعلية")

        if uploaded is None:
            st.warning("⬆️ يرجى رفع ملف Excel من الشريط الجانبي للبدء")
            st.stop()

        # ── تحميل البيانات ──
        @st.cache_data
        def load_excel(file):
            xl = pd.ExcelFile(file)
            return {{name: xl.parse(name) for name in xl.sheet_names}}

        sheets = load_excel(uploaded)
        sheet_name = st.selectbox("اختر الورقة", list(sheets.keys()))
        df = sheets[sheet_name].copy()

        # ── تنظيف تلقائي ──
        df.dropna(how="all", inplace=True)
        num_cols = df.select_dtypes(include=[np.number]).columns
        df[num_cols] = df[num_cols].fillna(df[num_cols].median())

        # ── KPIs ──
        st.subheader("📈 المؤشرات الرئيسية")
        cols = st.columns(4)
        numeric = df.select_dtypes(include=[np.number])
        for i, col in enumerate(list(numeric.columns)[:4]):
            with cols[i]:
                total = numeric[col].sum()
                mean  = numeric[col].mean()
                st.metric(col, f"{{total:,.0f}}", f"متوسط: {{mean:,.0f}}")

        st.markdown("---")

        # ── البيانات الخام ──
        with st.expander("📋 عرض البيانات المنظّفة"):
            st.dataframe(df, use_container_width=True)
            st.caption(f"{{len(df):,}} صف × {{len(df.columns)}} عمود")

        # ── الرسوم البيانية ──
        st.subheader("📊 الرسوم البيانية")
        cat_cols = [c for c in df.columns if df[c].nunique() < 20 and df[c].dtype == object]
        val_cols = list(numeric.columns)

        if cat_cols and val_cols:
            c1, c2 = st.columns(2)
            with c1:
                cat = st.selectbox("فئة المحور X", cat_cols, key="bar_cat")
                val = st.selectbox("قيمة المحور Y", val_cols, key="bar_val")
                grp = df.groupby(cat)[val].sum().sort_values(ascending=False).head(10)
                fig, ax = plt.subplots(figsize=(8,5))
                ax.bar(range(len(grp)), grp.values,
                       color=["#2E86C1","#F39C12","#27AE60","#E74C3C","#8E44AD"] * 3)
                ax.set_xticks(range(len(grp)))
                ax.set_xticklabels([str(x)[:15] for x in grp.index], rotation=30, ha="right")
                ax.set_title(f"{{cat}} vs {{val}}", fontweight="bold")
                ax.spines[["top","right"]].set_visible(False)
                st.pyplot(fig)
                plt.close()
            with c2:
                grp2 = df.groupby(cat)[val].sum().sort_values(ascending=False).head(6)
                fig2, ax2 = plt.subplots(figsize=(7,6))
                ax2.pie(grp2.values, labels=[str(x)[:12] for x in grp2.index],
                        autopct="%1.1f%%",
                        colors=["#2E86C1","#F39C12","#27AE60","#E74C3C","#8E44AD","#17A589"],
                        wedgeprops=dict(edgecolor="white", linewidth=1.5))
                ax2.set_title(f"توزيع {{cat}}", fontweight="bold")
                st.pyplot(fig2)
                plt.close()

        # ── مصفوفة الارتباط ──
        if len(val_cols) >= 2:
            st.subheader("🌡️ مصفوفة الارتباط")
            corr = numeric.corr()
            fig3, ax3 = plt.subplots(figsize=(10, max(4, len(corr)-1)))
            im = ax3.imshow(corr.values, cmap="RdYlGn", vmin=-1, vmax=1)
            ax3.set_xticks(range(len(corr.columns)))
            ax3.set_yticks(range(len(corr.columns)))
            ax3.set_xticklabels([c[:12] for c in corr.columns], rotation=45, ha="right")
            ax3.set_yticklabels([c[:12] for c in corr.columns])
            for i in range(len(corr)):
                for j in range(len(corr)):
                    ax3.text(j, i, f"{{corr.values[i,j]:.2f}}", ha="center", va="center", fontsize=8)
            plt.colorbar(im, ax=ax3)
            st.pyplot(fig3)
            plt.close()

        # ── الإحصاء الوصفي ──
        st.subheader("📐 الإحصاء الوصفي")
        st.dataframe(df.describe().round(2), use_container_width=True)

        st.markdown("---")
        st.caption("🤖 Smart Data Analyzer | تحليل ذكي آلي | {datetime.now().year}")
        ''')

        out_path = str(self.output_dir / "streamlit_dashboard.py")
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(code)
        print(f"✅  Streamlit Dashboard محفوظ: {out_path}")
        return out_path


# ══════════════════════════════════════════════
# MODULE 10 – Google Drive Integration (Template)
# ══════════════════════════════════════════════
GOOGLE_DRIVE_CODE = '''"""
Google Drive Integration – Smart Data Analyzer
===============================================
يراقب مجلداً محدداً في Google Drive ويُشغّل التحليل تلقائياً عند رصد ملف Excel جديد.

المتطلبات:
    pip install google-auth google-auth-oauthlib google-auth-httplib2 google-api-python-client

الإعداد:
1. أنشئ مشروعاً في Google Cloud Console
2. فعّل Google Drive API
3. حمّل ملف credentials.json
4. شغّل البرنامج وأدخل رمز المصادقة
"""

import os, time, pickle, tempfile
from pathlib import Path
from googleapiclient.discovery import build
from googleapiclient.http      import MediaIoBaseDownload
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request

SCOPES         = ["https://www.googleapis.com/auth/drive.readonly"]
WATCHED_FOLDER = "YOUR_FOLDER_ID"   # ← ضع معرّف المجلد هنا
POLL_INTERVAL  = 60                  # فحص كل 60 ثانية
PROCESSED_FILE = "processed_files.txt"


def get_drive_service():
    creds = None
    if os.path.exists("token.pickle"):
        with open("token.pickle", "rb") as f:
            creds = pickle.load(f)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file("credentials.json", SCOPES)
            creds = flow.run_local_server(port=0)
        with open("token.pickle", "wb") as f:
            pickle.dump(creds, f)
    return build("drive", "v3", credentials=creds)


def list_xlsx_files(service, folder_id):
    query = (f"'{folder_id}' in parents "
             f"and mimeType='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' "
             f"and trashed=false")
    results = service.files().list(q=query,
                                    fields="files(id,name,modifiedTime)",
                                    orderBy="modifiedTime desc").execute()
    return results.get("files", [])


def download_file(service, file_id, dest_path):
    request = service.files().get_media(fileId=file_id)
    with open(dest_path, "wb") as fh:
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            _, done = downloader.next_chunk()


def load_processed():
    if not os.path.exists(PROCESSED_FILE):
        return set()
    with open(PROCESSED_FILE) as f:
        return set(f.read().splitlines())


def save_processed(processed):
    with open(PROCESSED_FILE, "w") as f:
        f.write("\\n".join(processed))


def watch_and_analyze():
    from smart_data_analyzer import SmartDataAnalyzer
    service   = get_drive_service()
    processed = load_processed()
    print(f"👁️  مراقبة المجلد: {WATCHED_FOLDER}")
    print(f"⏱️  الفاصل الزمني: {POLL_INTERVAL} ثانية")

    while True:
        try:
            files = list_xlsx_files(service, WATCHED_FOLDER)
            for f in files:
                fid = f["id"]
                if fid in processed:
                    continue
                print(f"\\n🆕  ملف جديد: {f['name']}")
                with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
                    tmp_path = tmp.name
                download_file(service, fid, tmp_path)
                analyzer = SmartDataAnalyzer(tmp_path)
                analyzer.run()
                processed.add(fid)
                save_processed(processed)
                os.unlink(tmp_path)
        except Exception as e:
            print(f"⚠️  خطأ: {e}")
        time.sleep(POLL_INTERVAL)


if __name__ == "__main__":
    watch_and_analyze()
'''

EMAIL_CODE = '''"""
Email Notifier – Smart Data Analyzer
=====================================
يرسل التقارير النهائية تلقائياً عبر البريد الإلكتروني.

الإعداد:
    export EMAIL_SENDER="your@email.com"
    export EMAIL_PASSWORD="app_password"
    export EMAIL_RECIPIENT="recipient@email.com"
"""

import os, smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text      import MIMEText
from email.mime.base      import MIMEBase
from email               import encoders
from pathlib             import Path


def send_report(pdf_path: str, xlsx_path: str, pptx_path: str,
                subject: str = "تقرير تحليل البيانات – Smart Data Analyzer"):
    sender    = os.getenv("EMAIL_SENDER",    "sender@example.com")
    password  = os.getenv("EMAIL_PASSWORD",  "")
    recipient = os.getenv("EMAIL_RECIPIENT", "recipient@example.com")

    msg = MIMEMultipart()
    msg["From"]    = sender
    msg["To"]      = recipient
    msg["Subject"] = subject

    body = """
    السادة المحترمون،
    
    يُرفق طيّه تقرير التحليل الذكي للبيانات المُنشأ تلقائياً بواسطة نظام Smart Data Analyzer.
    
    المرفقات:
    • تقرير PDF احترافي
    • ملف Excel بالنتائج والتحليلات
    • عرض PowerPoint التقديمي
    
    مع التقدير،
    نظام Smart Data Analyzer
    """
    msg.attach(MIMEText(body, "plain", "utf-8"))

    for filepath in [pdf_path, xlsx_path, pptx_path]:
        if not filepath or not Path(filepath).exists():
            continue
        with open(filepath, "rb") as f:
            part = MIMEBase("application", "octet-stream")
            part.set_payload(f.read())
        encoders.encode_base64(part)
        part.add_header("Content-Disposition",
                        f"attachment; filename={Path(filepath).name}")
        msg.attach(part)

    with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
        server.login(sender, password)
        server.sendmail(sender, recipient, msg.as_string())
    print(f"✅  تم إرسال التقرير إلى: {recipient}")


if __name__ == "__main__":
    import glob
    outputs = Path("outputs")
    pdfs  = sorted(glob.glob(str(outputs / "report_*.pdf")))
    xlsxs = sorted(glob.glob(str(outputs / "results_*.xlsx")))
    pptxs = sorted(glob.glob(str(outputs / "presentation_*.pptx")))
    if pdfs:
        send_report(pdfs[-1],
                    xlsxs[-1] if xlsxs else "",
                    pptxs[-1] if pptxs else "")
'''


# ══════════════════════════════════════════════
# MODULE 11 – المنسق الرئيسي
# ══════════════════════════════════════════════
class SmartDataAnalyzer:
    """ينسّق جميع مراحل التحليل ويُنتج المخرجات الكاملة"""

    def __init__(self, filepath: str):
        self.filepath = filepath

    def run(self) -> dict:
        print("\n" + "═"*60)
        print("  📊  Smart Data Analyzer – بدء التحليل")
        print("═"*60 + "\n")

        # ── 1. تحميل وتنظيف ──
        loader = DataLoader(self.filepath)
        loader.load().clean_data()

        classifier = ColumnClassifier()
        analyzer   = StatisticalAnalyzer()
        charts_gen = ChartGenerator(OUTPUT_DIR)

        all_analysis = {}
        all_charts   = []

        # ── 2. التحليل والرسوم لكل ورقة ──
        for sheet_name, df in loader.clean.items():
            print(f"\n📋  تحليل الورقة: {sheet_name}")
            col_map  = classifier.classify(df)
            analysis = analyzer.analyze(df, col_map)
            all_analysis[sheet_name] = analysis

            # الرسوم البيانية
            if col_map["category"] and col_map["value"]:
                c, v = col_map["category"][0], col_map["value"][0]
                all_charts.append(
                    charts_gen.bar_chart(df, c, v, f"مخطط الأعمدة – {sheet_name}"))
                all_charts.append(
                    charts_gen.pie_chart(df, c, v, f"التوزيع النسبي – {sheet_name}"))

            if "trend" in analysis:
                all_charts.append(
                    charts_gen.line_chart(analysis["trend"]["monthly"],
                                          col_map["value"][0] if col_map["value"] else "القيمة",
                                          f"الاتجاه الزمني – {sheet_name}"))

            if col_map["value"] and len(col_map["value"]) >= 2:
                all_charts.append(
                    charts_gen.scatter_plot(df,
                                            col_map["value"][0], col_map["value"][1],
                                            f"العلاقة – {sheet_name}"))

            hm = charts_gen.heatmap(df, f"مصفوفة الارتباط – {sheet_name}")
            if hm:
                all_charts.append(hm)

            dash = charts_gen.dashboard(df, col_map, analysis, sheet_name)
            if dash:
                all_charts.append(dash)

        # حذف المسارات الفارغة
        all_charts = [c for c in all_charts if c]

        # ── 3. التقرير العربي ──
        print("\n✍️   كتابة التقرير التفسيري...")
        writer      = ArabicReportWriter()
        report_text = writer.write(loader, all_analysis)
        report_txt_path = str(OUTPUT_DIR / f"report_{TIMESTAMP}.txt")
        with open(report_txt_path, "w", encoding="utf-8") as f:
            f.write(report_text)

        # ── 4. PDF ──
        print("📄  إنشاء تقرير PDF...")
        pdf_gen  = PDFReportGenerator(OUTPUT_DIR)
        pdf_path = pdf_gen.generate(report_text, all_charts, list(loader.clean.keys())[0])

        # ── 5. Excel ──
        print("📊  إنشاء ملف Excel بالنتائج...")
        xlsx_gen  = ExcelResultsGenerator(OUTPUT_DIR)
        xlsx_path = xlsx_gen.generate(loader, all_analysis)

        # ── 6. PowerPoint ──
        print("🎨  إنشاء عرض PowerPoint...")
        pptx_gen  = PowerPointGenerator(OUTPUT_DIR)
        pptx_path = pptx_gen.generate(loader, all_analysis, all_charts)

        # ── 7. Streamlit ──
        print("🌐  إنشاء Streamlit Dashboard...")
        st_gen  = StreamlitDashboardGenerator(OUTPUT_DIR)
        st_path = st_gen.generate(self.filepath)

        # ── 8. Google Drive + Email Code ──
        gdrive_path = str(OUTPUT_DIR / "google_drive_watcher.py")
        email_path  = str(OUTPUT_DIR / "email_notifier.py")
        with open(gdrive_path, "w", encoding="utf-8") as f:
            f.write(GOOGLE_DRIVE_CODE)
        with open(email_path, "w", encoding="utf-8") as f:
            f.write(EMAIL_CODE)

        # ── ملخص ──
        print("\n" + "═"*60)
        print("  ✅  اكتمل التحليل! المخرجات:")
        print("═"*60)
        outputs = {
            "pdf":        pdf_path,
            "excel":      xlsx_path,
            "powerpoint": pptx_path,
            "streamlit":  st_path,
            "gdrive":     gdrive_path,
            "email":      email_path,
            "report_txt": report_txt_path,
            "charts":     all_charts,
        }
        for k, v in outputs.items():
            if isinstance(v, list):
                print(f"  📊  رسوم بيانية : {len(v)} ملف")
            else:
                print(f"  📁  {k:12}: {Path(v).name}")
        print("═"*60 + "\n")
        return outputs


# ══════════════════════════════════════════════
# نقطة الدخول
# ══════════════════════════════════════════════
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("الاستخدام: python smart_data_analyzer.py <path_to_excel.xlsx>")
        print("مثال:     python smart_data_analyzer.py sales_data.xlsx")
        sys.exit(1)

    filepath = sys.argv[1]
    if not os.path.exists(filepath):
        print(f"❌  الملف غير موجود: {filepath}")
        sys.exit(1)

    analyzer = SmartDataAnalyzer(filepath)
    analyzer.run()
